"""Helpers for building a deck on top of the CYBERFORT pilot pptx template.

The template ships four slides:
  0  Title Slide       — branded cover (background photo, CRA subtitle, presenter block)
  1  Title and Content  — branded content slide (logo top-right, EU/partner logos bottom)
  2  Title and Content  — second copy of the same branded content slide
  3  Title Slide        — "Thank You" closer

Content slides are produced by deep-copying slide 1 so every clone keeps the
branding shapes, then rewriting the title and drawing into the body area.

Safe content band on a 13.33 x 7.50in slide, given the template's branding:
  x: 0.64 .. 12.70   (width 12.06)
  y: 1.85 .. 6.35    (top-right logo ends at 1.88; bottom logo strip starts 6.44)
"""

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
R_ID = _R_NS + "id"
# every attribute in DrawingML that carries a relationship id
REL_ATTRS = (_R_NS + "embed", _R_NS + "link", R_ID)

# Theme "Dividend"
NAVY = RGBColor(0x1A, 0x32, 0x60)
BLUE = RGBColor(0x45, 0x90, 0xB8)
CYAN = RGBColor(0x45, 0xCB, 0xE8)
GREY = RGBColor(0x96, 0x9F, 0xA7)
OLIVE = RGBColor(0xA2, 0xC7, 0x77)
GREEN = RGBColor(0x2E, 0x7D, 0x4F)
AMBER = RGBColor(0xB8, 0x7A, 0x12)
RED = RGBColor(0xA8, 0x32, 0x32)
INK = RGBColor(0x3D, 0x3D, 0x3D)
MUTED = RGBColor(0x6B, 0x72, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF2, 0xF5, 0xF9)
WASH_G = RGBColor(0xEC, 0xF5, 0xEF)
WASH_R = RGBColor(0xFB, 0xEF, 0xEF)
WASH_A = RGBColor(0xFD, 0xF7, 0xE8)
LINE = RGBColor(0xD8, 0xDE, 0xE6)

MAJOR = "Gill Sans MT"
MINOR = "Corbel"

PROTOTYPE = 1
BAND_L, BAND_R = 0.64, 12.70
BAND_W = BAND_R - BAND_L

# The layout draws a navy header rectangle from y=0.67 to y=1.97 and the footer
# logo strip starts at y=6.44. The header holds the title and its standfirst and
# nothing else; content lives below it with a deliberate white gutter between.
HEADER_B = 1.97      # bottom of the navy header rectangle
TITLE_T = 0.72       # title text, top-anchored so long titles grow downward
TITLE_H = 0.80       # room for two lines at 26pt
KICK_Y = 1.58        # standfirst, still inside the header
TOP = 2.22           # first content row — 0.25in of white below the header
BOT = 6.40           # last content pixel — the footer strip begins at 6.44


# ---------------------------------------------------------------- slide plumbing

class Pristine:
    """A snapshot of an untouched prototype slide, used to stamp out clones.

    Cloning from the live prototype is wrong once anything has been drawn on it —
    every clone would inherit the earlier slide's content. Snapshot the shape XML
    before filling anything and stamp from the snapshot instead.
    """

    def __init__(self, prs, src_index=PROTOTYPE):
        src = prs.slides[src_index]
        self.layout = src.slide_layout
        self.part = src.part
        self.shapes = [copy.deepcopy(el) for el in list(src.shapes._spTree)
                       if el.tag.endswith(('}sp', '}pic', '}graphicFrame',
                                           '}grpSp', '}cxnSp'))]

    def stamp(self, prs):
        """Append a fresh copy of the prototype slide.

        The copied shape XML addresses its images by literal rId — "Picture 7"
        embeds rId3. The new slide part hands out its own rIds in its own order,
        so the copied references must be remapped or the images silently swap
        (which is how the CYBERFORT logo turned into an EU flag from slide 4 on).
        """
        dst = prs.slides.add_slide(self.layout)
        for shp in list(dst.shapes):
            shp._element.getparent().remove(shp._element)

        remap = {}
        for rId, rel in self.part.rels.items():
            if rel.reltype == RT.SLIDE_LAYOUT:
                continue                      # add_slide() already made this one
            remap[rId] = (dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                          if rel.is_external
                          else dst.part.rels.get_or_add(rel.reltype, rel.target_part))

        for el in self.shapes:
            clone = copy.deepcopy(el)
            for node in clone.iter():
                for attr in REL_ATTRS:
                    old = node.get(attr)
                    if old in remap:
                        node.set(attr, remap[old])
            dst.shapes._spTree.append(clone)
        return dst


def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    prs.part.drop_rel(sldId.get(R_ID))
    sldIdLst.remove(sldId)


def reorder(prs, slides):
    """Reorder the deck to exactly the given slide objects, in order."""
    sldIdLst = prs.slides._sldIdLst
    by_part = {}
    for sldId in list(sldIdLst):
        by_part[prs.part.related_part(sldId.get(R_ID))] = sldId
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for s in slides:
        sldIdLst.append(by_part[s.part])


def placeholder(slide, idx):
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == idx:
            return shp
    return None


def set_title(slide, text, size=26, width=None):
    """Title sits on the layout's navy header band, so it must stay white.

    The placeholder ships bottom-anchored, which makes a two-line title grow
    upwards into the standfirst. Pin it top-anchored inside the header instead.
    """
    ph = placeholder(slide, 0)
    if ph is None:
        return None
    if width:
        ph.width = Inches(width)
    ph.top = Inches(TITLE_T)
    ph.height = Inches(TITLE_H)
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.line_spacing = 0.92
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = MAJOR
    r.font.color.rgb = WHITE
    return ph


def drop_body(slide):
    ph = placeholder(slide, 1)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def slide_number(slide, l=12.42, t=6.68, w=0.80, h=0.32, size=10, color=GREY):
    """Live slide-number field, parked right of the footer logo strip."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "#"
    run.font.size = Pt(size)
    run.font.name = MINOR
    run.font.bold = True
    run.font.color.rgb = color
    # swap the run for a slidenum field so PowerPoint renders the real number
    rPr = run._r.get_or_add_rPr()
    fld = parse_xml(
        '<a:fld %s id="{9C2E5A31-7B4D-4E68-9F13-6D4A1C0B27E5}" type="slidenum">'
        "%s<a:t>#</a:t></a:fld>" % (nsdecls("a"), rPr.xml))
    run._r.addprevious(fld)
    run._r.getparent().remove(run._r)
    return box


def content_slide(prs, title, title_size=26, title_width=8.3, proto=None,
                  number=True):
    slide = proto
    set_title(slide, title, size=title_size, width=title_width)
    drop_body(slide)
    if number:
        slide_number(slide)
    return slide


# ---------------------------------------------------------------- drawing

def text(slide, l, t, w, h, lines, size=13, color=INK, font=MINOR, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.0, space_after=4, anchor=MSO_ANCHOR.TOP):
    """lines: str | list of str | list of (text, dict-of-overrides)."""
    if isinstance(lines, str):
        lines = [lines]
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    first = True
    for item in lines:
        opts = {}
        if isinstance(item, tuple):
            item, opts = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("spacing", spacing)
        p.space_after = Pt(opts.get("space_after", space_after))
        r = p.add_run()
        r.text = item
        r.font.size = Pt(opts.get("size", size))
        r.font.name = opts.get("font", font)
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
    return box


def card(slide, l, t, w, h, fill=WASH, line=None, line_w=0.75, radius=0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def bar(slide, l, t, w, h, fill, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                  Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def accent_rule(slide, l, t, w=1.5, h=0.055, fill=CYAN):
    return bar(slide, l, t, w, h, fill)


def kicker(slide, l, t=None, label="", color=CYAN, size=10.5):
    """Standfirst under the title. Still inside the navy header — keep it cyan."""
    box = text(slide, l, KICK_Y if t is None else t, 8.3, 0.26, label.upper(),
               size=size, color=color, bold=True, font=MINOR, spacing=1.0)
    box.name = "Standfirst"          # part of the header block, never nudged
    return box


BRAND_SHAPES = {"Group 3", "Picture 7", "Picture 8", "Picture 9", "Standfirst"}


def is_chrome(shape):
    """True for template branding, the title block, or the slide-number field."""
    if shape.is_placeholder and shape.placeholder_format.idx == 0:
        return True
    if shape.name in BRAND_SHAPES:
        return True
    return "slidenum" in shape._element.xml


def nudge_content(prs, dy, top=TOP, bot=BOT, report=True):
    """Push every slide's content down by dy so it clears the header band.

    Applied once, after all slides are drawn, so individual slide code can stay
    written against a single origin. The shift is clamped per slide so nothing is
    pushed into the footer logo strip; a clamped slide is reported, not silently
    squeezed.
    """
    for i, slide in enumerate(prs.slides, 1):
        if slide.slide_layout.name == "Title Slide":
            continue                      # the cover and the closer are full-bleed
        content = [sh for sh in slide.shapes
                   if sh.top is not None and not is_chrome(sh)]
        if not content:
            continue
        y0 = min(Emu(sh.top).inches for sh in content)
        y1 = max(Emu(sh.top).inches + Emu(sh.height).inches for sh in content)
        if y0 > top:                      # already clear (e.g. the cover)
            continue
        eff = min(dy, bot - y1)
        if eff <= 0.005:
            if report:
                print(f"  !! slide {i}: no room to shift (content ends {y1:.2f})")
            continue
        for sh in content:
            sh.top = Emu(sh.top) + Inches(eff)
        if report and eff < dy - 0.005:
            print(f"  ~  slide {i}: shifted {eff:.2f}in of {dy:.2f} "
                  f"(content now {y0 + eff:.2f}-{y1 + eff:.2f})")
    return prs


def stat_tile(slide, l, t, w, h, value, label, note=None,
              value_color=NAVY, fill=WASH, line=LINE, value_size=34):
    card(slide, l, t, w, h, fill=fill, line=line)
    text(slide, l + 0.18, t + 0.16, w - 0.36, 0.62, value,
         size=value_size, color=value_color, bold=True, font=MAJOR, spacing=0.9)
    text(slide, l + 0.18, t + 0.16 + value_size / 62.0, w - 0.36, 0.5, label,
         size=11.5, color=INK, bold=True)
    if note:
        text(slide, l + 0.18, t + h - 0.52, w - 0.36, 0.44, note,
             size=9.5, color=MUTED, spacing=0.95)


def picture(slide, path, l, t, w):
    return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def framed_picture(slide, path, l, t, w, pad=0.06, caption=None):
    """Place a screenshot with a hairline frame behind it."""
    pic = picture(slide, path, l, t, w)
    h = pic.height / 914400
    frame = card(slide, l - pad, t - pad, w + 2 * pad, h + 2 * pad,
                 fill=WHITE, line=LINE, line_w=1.0, radius=0.02)
    spTree = slide.shapes._spTree
    spTree.remove(frame._element)
    spTree.insert(list(spTree).index(pic._element), frame._element)
    if caption:
        text(slide, l, t + h + pad + 0.08, w, 0.3, caption, size=9.5,
             color=MUTED, align=PP_ALIGN.CENTER)
    return pic


def retext(shape, lines):
    """Replace a shape's paragraph texts in place, keeping each run's formatting."""
    paras = shape.text_frame.paragraphs
    for i, new in enumerate(lines):
        if i >= len(paras):
            break
        p = paras[i]
        runs = [r for r in p.runs]
        if not runs:
            continue
        runs[0].text = new
        for extra in runs[1:]:
            extra.text = ""
